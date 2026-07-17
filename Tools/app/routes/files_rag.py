import base64
import json
import secrets
import re
import os
import uuid
import datetime
import io
from io import BytesIO
import pdfplumber
from fastapi import APIRouter, Depends, HTTPException, Response, Request
from fastapi.responses import JSONResponse, RedirectResponse
from flask_jwt_extended import decode_token
from pydantic import BaseModel
from pgvector.sqlalchemy import Vector
from sqlalchemy import Column, String, Text, DateTime, JSON
from sqlalchemy.orm import Session
from sqlalchemy.dialects.postgresql import UUID
from googleapiclient.discovery import build
from googleapiclient.http import MediaIoBaseDownload
from google.oauth2.credentials import Credentials
from google.auth.transport.requests import Request as GoogleRequest
from google_auth_oauthlib.flow import Flow
from app.BrightpathAI.models.users import User
from app.db.database import Base, get_db
import dotenv
from docx import Document
from pptx import Presentation
import openpyxl
from PIL import Image
import pytesseract
import mimetypes
from striprtf.striprtf import rtf_to_text


dotenv.load_dotenv()
router = APIRouter(prefix="/files", tags=["Drive Ingest"])
GOOGLE_SCOPES = ["https://www.googleapis.com/auth/drive.readonly"]
FOLDER_MIME = "application/vnd.google-apps.folder"
EXPORT_MAP = {
    "application/vnd.google-apps.document": "text/plain",
    "application/vnd.google-apps.spreadsheet": "text/csv",
    "application/vnd.google-apps.presentation": "text/plain",
}

class File(Base):
    __tablename__ = "files"

    id = Column(UUID(as_uuid=True), primary_key=True, default=uuid.uuid4)
    source = Column(String)
    user_email = Column(String)
    drive_file_id = Column(String, index=True)
    file_name = Column(String)
    file_path = Column(Text)
    mime_type = Column(String)
    content = Column(Text)
    created_at = Column(DateTime, default=datetime.datetime.utcnow)

class DriveToken(Base):
    __tablename__ = "drive_tokens"

    id = Column(UUID(as_uuid=True), primary_key=True, default=uuid.uuid4)
    user_email = Column(String, unique=True, index=True)
    refresh_token = Column(Text, nullable=False)
    created_at = Column(DateTime, default=datetime.datetime.utcnow)

class IngestRequest(BaseModel):
    drive_link: str
    user_email: str

def build_flow(state: str):
    return Flow.from_client_secrets_file(
        os.getenv("GOOGLE_CLIENT_SECRET_FILE"),
        scopes=GOOGLE_SCOPES,
        redirect_uri=os.getenv("GOOGLE_REDIRECT_URI"),
        state=state,
    )

def get_current_user(request: Request, db: Session = Depends(get_db)):
    """
    Extract user from JWT in cookie
    """
    token = request.cookies.get("access_token")
    if not token:
        raise HTTPException(status_code=401, detail="Not authenticated")
    payload = decode_token(token)
    user = db.query(User).filter(User.id == payload["sub"]).first()
    if not user:
        raise HTTPException(status_code=401, detail="Invalid user")
    return user

def get_drive_service(refresh_token: str):
    creds = Credentials(
        token=None,
        refresh_token=refresh_token,
        token_uri="https://oauth2.googleapis.com/token",
        client_id=os.getenv("GOOGLE_CLIENT_ID"),
        client_secret=os.getenv("GOOGLE_CLIENT_SECRET"),
        scopes=GOOGLE_SCOPES,
    )
    creds.refresh(GoogleRequest())
    return build("drive", "v3", credentials=creds)

def parse_drive_link(link: str):
    if "/folders/" in link:
        return "folder", re.search(r"/folders/([^/?]+)", link).group(1)
    if "/d/" in link:
        return "file", re.search(r"/d/([^/?]+)", link).group(1)
    return "root", "root"

def list_recursive(service, folder_id, path=""):
    results = []
    page_token = None
    while True:
        res = service.files().list(
            q=f"'{folder_id}' in parents and trashed=false",
            fields="nextPageToken, files(id,name,mimeType,owners(emailAddress))",
            pageToken=page_token
        ).execute()
        for f in res.get("files", []):
            full_path = f"{path}/{f['name']}"
            f["path"] = full_path
            if f["mimeType"] == FOLDER_MIME:
                results.extend(list_recursive(service, f["id"], full_path))
            else:
                results.append(f)
        page_token = res.get("nextPageToken")
        if not page_token:
            break
    return results

def get_single_file(service, file_id):
    f = service.files().get(
        fileId=file_id,
        fields="id,name,mimeType,owners(emailAddress)"
    ).execute()
    f["path"] = f["name"]
    return [f]

def extract_text_from_bytes(data: bytes, filename: str) -> str | None:
    ext = filename.lower().split(".")[-1]
    try:
        if ext in {"txt", "csv", "json", "md"}:
            return data.decode("utf-8", errors="ignore").strip() or None
        if ext == "pdf":
            text = ""
            with pdfplumber.open(io.BytesIO(data)) as pdf:
                for p in pdf.pages:
                    t = p.extract_text()
                    if t:
                        text += t + "\n"
            return text.strip() or None
        if ext == "docx":
            doc = Document(io.BytesIO(data))
            return "\n".join(p.text for p in doc.paragraphs).strip() or None
        if ext == "pptx":
            prs = Presentation(io.BytesIO(data))
            text = ""
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
            return text.strip() or None
        if ext == "xlsx":
            wb = openpyxl.load_workbook(io.BytesIO(data), data_only=True)
            text = ""
            for sheet in wb.worksheets:
                for row in sheet.iter_rows(values_only=True):
                    text += " ".join(str(c) for c in row if c) + "\n"
            return text.strip() or None
        if ext == "rtf":
            return rtf_to_text(data.decode("utf-8", errors="ignore")).strip() or None
        if ext in {"png", "jpg", "jpeg", "bmp", "tiff"}:
            img = Image.open(io.BytesIO(data))
            return pytesseract.image_to_string(img).strip() or None
    except Exception:
        return None
    return None

def download_and_extract(service, file):
    file_id = file["id"]
    mime_type = file["mimeType"]
    name = file["name"]
    try:
        if mime_type.startswith("application/vnd.google-apps"):
            export_mime = EXPORT_MAP.get(mime_type, "text/plain")
            request = service.files().export_media(
                fileId=file_id,
                mimeType=export_mime
            )
        else:
            request = service.files().get_media(fileId=file_id)
        fh = io.BytesIO()
        downloader = MediaIoBaseDownload(fh, request)
        done = False
        while not done:
            _, done = downloader.next_chunk()
        data = fh.getvalue()
        if not data:
            return None
        return extract_text_from_bytes(data, name)
    except Exception as e:
        print(f"Failed to extract {name}: {e}")
        return None

@router.get("/oauth/start")
def oauth_start(
    response: Response,
    email: str,
    request: Request,
    db: Session = Depends(get_db),
):
    """
    Start OAuth flow and embed user's email in state
    """
    state_data = {"user_email": email, "nonce": secrets.token_urlsafe(8)}
    state_encoded = base64.urlsafe_b64encode(json.dumps(state_data).encode()).decode()
    flow = build_flow(state_encoded)
    auth_url, _ = flow.authorization_url(
        access_type="offline",
        prompt="consent",
        include_granted_scopes="true",
    )
    response.set_cookie(
        key="oauth_state",
        value=state_encoded,
        httponly=True,
        samesite="lax",
        path="/",
    )
    return {"oauth_url": auth_url}

@router.get("/oauth/callback")
def oauth_callback(
    request: Request,
    code: str,
    state: str,
    db: Session = Depends(get_db)
):
    """
    Callback receives Google auth code. User email is extracted from state.
    """
    saved_state = request.cookies.get("oauth_state")
    if not saved_state:
        raise HTTPException(status_code=400, detail="Missing OAuth state")

    def decode_state(s: str) -> dict:
        return json.loads(base64.urlsafe_b64decode(s + "==").decode())

    try:
        cookie_state = decode_state(saved_state)
        query_state = decode_state(state)
    except Exception:
        raise HTTPException(status_code=400, detail="Invalid OAuth state encoding")
    if cookie_state["nonce"] != query_state["nonce"]:
        raise HTTPException(status_code=400, detail="OAuth state mismatch")
    email = query_state["user_email"]

    flow = build_flow(state)
    flow.fetch_token(code=code)
    creds = flow.credentials

    token = db.query(DriveToken).filter_by(user_email=email).first()
    if token:
        if creds.refresh_token:
            token.refresh_token = creds.refresh_token
    else:
        db.add(DriveToken(user_email=email, refresh_token=creds.refresh_token))
    db.commit()
    return {"status": "OAuth success", "email": email}

@router.post("/ingest/public-drive")
def ingest(payload: IngestRequest, db: Session = Depends(get_db)):
    token = db.query(DriveToken).filter_by(user_email=payload.user_email).first()
    if not token:
        return JSONResponse(
            {"error": "Login required", "oauth": "/files/oauth/start"},
            status_code=401
        )
    service = get_drive_service(token.refresh_token)
    target_type, target_id = parse_drive_link(payload.drive_link)
    if target_type == "root":
        files = list_recursive(service, "root")
    elif target_type == "folder":
        files = list_recursive(service, target_id)
    else:
        files = get_single_file(service, target_id)
    for f in files:
        customer_email = payload.user_email
        owners = f.get("owners", [])
        modifiedTime = datetime.fromisoformat(
            f["modifiedTime"].replace("Z", "+00:00")
        )   
        if owners and owners[0].get("emailAddress"):
            customer_email = owners[0]["emailAddress"]
        exists = db.query(File).filter(File.drive_file_id == f["id"], File.created_at >= modifiedTime, File.user_email == customer_email).first()
        if exists:
            continue
        text = download_and_extract(service, f)
        if not text or not text.strip():
            continue
        db.add(File(
            source="public_drive",
            user_email=customer_email,
            drive_file_id=f["id"],
            file_name=f["name"],
            file_path=f["path"],
            mime_type=f["mimeType"],
            content=text
        ))
    db.commit()
    return {"status": "ingestion complete"}
